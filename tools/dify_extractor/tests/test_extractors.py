from io import BytesIO

import pytest
import xlwt
from docx import Document as WordDocument
from dify_plugin.invocations.file import UploadFileResponse
from openpyxl import Workbook
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from tools.context import ExtractionContext
from tools.csv_extractor import CSVExtractor
from tools.errors import ExtractionError
from tools.excel_extractor import ExcelExtractor
from tools.html_extractor import HtmlExtractor
from tools.image_assets import ImageAsset
from tools.json_extractor import JSONExtractor
from tools.markdown_extractor import MarkdownExtractor
from tools.pdf_extractor import PdfExtractor
from tools.pptx_extractor import PPTXExtractor
from tools.text_extractor import TextExtractor
from tools.word_extractor import WordExtractor
from tools.yaml_extractor import YAMLExtractor


def _context(
    blob: bytes,
    filename: str,
    extension: str,
    *,
    image_service=None,
) -> ExtractionContext:
    return ExtractionContext(
        file_bytes=blob,
        file_name=filename,
        file_extension=extension,
        image_service=image_service,
    )


def _save_workbook(workbook) -> bytes:
    buffer = BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def _minimal_pdf(text: str) -> bytes:
    stream = f"BT /F1 12 Tf 72 720 Td ({text}) Tj ET".encode()
    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        (
            b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
            b"/Resources << /Font << /F1 4 0 R >> >> /Contents 5 0 R >>"
        ),
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
        f"<< /Length {len(stream)} >>\nstream\n".encode() + stream + b"\nendstream",
    ]
    pdf = bytearray(b"%PDF-1.4\n")
    offsets = [0]
    for object_number, body in enumerate(objects, start=1):
        offsets.append(len(pdf))
        pdf.extend(f"{object_number} 0 obj\n".encode())
        pdf.extend(body)
        pdf.extend(b"\nendobj\n")
    xref_offset = len(pdf)
    pdf.extend(f"xref\n0 {len(objects) + 1}\n".encode())
    pdf.extend(b"0000000000 65535 f \n")
    for offset in offsets[1:]:
        pdf.extend(f"{offset:010d} 00000 n \n".encode())
    pdf.extend(
        (
            f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\n"
            f"startxref\n{xref_offset}\n%%EOF\n"
        ).encode()
    )
    return bytes(pdf)


def test_text_extractor_decodes_non_utf8():
    result = TextExtractor(_context("café".encode("cp1252"), "note.txt", ".txt")).extract()

    assert result.md_content == "café"
    assert result.documents[0].metadata == {"source": "note.txt"}


def test_html_extractor_removes_non_content_elements():
    html = b"<html><style>hidden</style><script>secret</script><body><h1>Hello</h1></body></html>"
    result = HtmlExtractor(_context(html, "page.html", ".html")).extract()

    assert result.md_content == "Hello"
    assert "hidden" not in result.md_content
    assert "secret" not in result.md_content


def test_json_extractor_formats_valid_json_and_rejects_invalid_json():
    valid = JSONExtractor(_context(b'{"city":"Shanghai"}', "data.json", ".json")).extract()

    assert '"city": "Shanghai"' in valid.md_content
    with pytest.raises(ExtractionError, match="line 1, column 2"):
        JSONExtractor(_context(b"{broken", "bad.json", ".json")).extract()


def test_yaml_extractor_formats_valid_yaml_and_rejects_invalid_yaml():
    valid = YAMLExtractor(_context("城市: 上海\n".encode(), "data.yaml", ".yaml")).extract()

    assert "城市: 上海" in valid.md_content
    assert valid.md_content.endswith("```")
    with pytest.raises(ExtractionError, match="invalid YAML"):
        YAMLExtractor(_context(b"items: [", "bad.yaml", ".yaml")).extract()
    with pytest.raises(ExtractionError, match="no YAML content"):
        YAMLExtractor(_context(b"# only a comment\n", "empty.yaml", ".yaml")).extract()


def test_csv_extractor_preserves_multiline_values_and_escapes_markdown():
    blob = b'name,note\nAda,"first|second\nthird"\n'
    result = CSVExtractor(_context(blob, "people.csv", ".csv")).extract()

    assert "first\\|second<br>third" in result.md_content
    assert result.documents[0].metadata == {"source": "people.csv", "row": 0}
    assert "first|second\nthird" in result.documents[0].page_content


def test_csv_extractor_rejects_inconsistent_column_count():
    with pytest.raises(ExtractionError, match="expected 2"):
        CSVExtractor(_context(b"a,b\n1\n", "bad.csv", ".csv")).extract()


def test_markdown_extractor_splits_sections_and_ignores_headings_in_fences():
    markdown = b"preamble\n# First\nbody\n```md\n# not a heading\n```\n## Second\nend"
    result = MarkdownExtractor(_context(markdown, "readme.md", ".md")).extract()

    assert [document.page_content for document in result.documents] == [
        "preamble",
        "First\nbody\n```md\n# not a heading\n```",
        "Second\nend",
    ]
    assert all(document.metadata == {"source": "readme.md"} for document in result.documents)


def test_markdown_remote_image_failure_preserves_original_link():
    class FailedImageService:
        def download_and_upload(self, _url):
            return None

    markdown = b"![diagram](https://example.com/diagram.png)"
    result = MarkdownExtractor(
        _context(markdown, "readme.md", ".md", image_service=FailedImageService())
    ).extract()

    assert result.md_content == markdown.decode()
    assert result.img_list == []


def test_markdown_remote_image_success_rewrites_link_and_returns_image():
    uploaded_file = UploadFileResponse(
        id="image-1",
        name="diagram.png",
        size=5,
        extension="png",
        mime_type="image/png",
        preview_url="https://dify.example/diagram.png",
    )

    class SuccessfulImageService:
        def download_and_upload(self, _url):
            return ImageAsset(
                file=uploaded_file,
                markdown="![image](https://dify.example/diagram.png)",
            )

    result = MarkdownExtractor(
        _context(
            b"![diagram](https://example.com/diagram.png)",
            "readme.md",
            ".md",
            image_service=SuccessfulImageService(),
        )
    ).extract()

    assert result.md_content == "![diagram](https://dify.example/diagram.png)"
    assert result.img_list == [uploaded_file]


def test_xlsx_extractor_preserves_hyperlinks_and_metadata():
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "People"
    sheet.append(["name", "profile"])
    sheet.append(["Ada", "site|home"])
    sheet["B2"].hyperlink = "https://example.com/ada"

    result = ExcelExtractor(_context(_save_workbook(workbook), "people.xlsx", ".xlsx")).extract()

    assert "[site\\|home](https://example.com/ada)" in result.md_content
    assert "[site|home](https://example.com/ada)" in result.documents[0].page_content
    assert result.documents[0].metadata == {
        "source": "people.xlsx",
        "sheet": "People",
        "row": 0,
    }


def test_xls_extractor_uses_native_xlrd_path():
    workbook = xlwt.Workbook()
    sheet = workbook.add_sheet("People")
    sheet.write(0, 0, "name")
    sheet.write(1, 0, "Ada")
    blob = _save_workbook(workbook)

    result = ExcelExtractor(_context(blob, "people.xls", ".xls")).extract()

    assert "## People" in result.md_content
    assert result.documents[0].page_content == "name: Ada"


def test_docx_extractor_handles_text_and_escaped_tables():
    document = WordDocument()
    paragraph = document.add_paragraph()
    paragraph.add_run("Hello ")
    paragraph.add_run("DOCX")
    table = document.add_table(rows=2, cols=1)
    table.cell(0, 0).text = "heading"
    table.cell(1, 0).text = "a|b"

    result = WordExtractor(_context(_save_workbook(document), "report.docx", ".docx")).extract()

    assert "Hello DOCX" in result.md_content
    assert "a\\|b" in result.md_content
    assert result.documents[0].metadata == {"source": "report.docx"}


def test_docx_embedded_image_failure_does_not_discard_text():
    class FailedImageService:
        def upload_embedded(self, *_args, **_kwargs):
            return None

    image = BytesIO()
    Image.new("RGB", (1, 1), "white").save(image, format="PNG")
    image.seek(0)
    document = WordDocument()
    document.add_paragraph("Text survives")
    document.add_picture(image)

    result = WordExtractor(
        _context(
            _save_workbook(document),
            "report.docx",
            ".docx",
            image_service=FailedImageService(),
        )
    ).extract()

    assert "Text survives" in result.md_content
    assert result.img_list == []


def test_pptx_extractor_supports_http_hyperlinks_and_escaped_tables():
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = textbox.text_frame.paragraphs[0].add_run()
    run.text = "Example"
    run.hyperlink.address = "http://example.com"
    table = slide.shapes.add_table(2, 1, Inches(1), Inches(2), Inches(4), Inches(1)).table
    table.cell(0, 0).text = "heading"
    table.cell(1, 0).text = "a|b"

    result = PPTXExtractor(_context(_save_workbook(presentation), "slides.pptx", ".pptx")).extract()

    assert "[Example](http://example.com)" in result.md_content
    assert "a\\|b" in result.md_content


def test_pdf_extractor_returns_page_documents_and_closes_resources():
    result = PdfExtractor(_context(_minimal_pdf("Hello PDF"), "report.pdf", ".pdf")).extract()

    assert "Hello PDF" in result.md_content
    assert result.documents[0].metadata == {"source": "report.pdf", "page": 0}


def test_pdf_extractor_rejects_invalid_pdf():
    with pytest.raises(ExtractionError, match="not a valid PDF"):
        PdfExtractor(_context(b"not a PDF", "bad.pdf", ".pdf")).extract()
