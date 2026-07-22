"""PPTX document extractor."""

from io import BytesIO
from urllib.parse import urlparse

from pptx import Presentation

from tools.document import Document, ExtractorResult
from tools.extractor_base import BaseExtractor
from tools.helpers import render_markdown_table


class PPTXExtractor(BaseExtractor):
    def extract(self) -> ExtractorResult:
        content, image_files = self._parse_pptx()
        return ExtractorResult(
            md_content=content,
            documents=[
                Document(
                    page_content=content,
                    metadata={"source": self.context.file_name},
                )
            ],
            img_list=image_files,
        )

    def _extract_images(self, presentation):
        image_map = {}
        image_files = []
        image_service = self.context.image_service
        if image_service is None:
            return image_map, image_files

        for slide_index, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                image = getattr(shape, "image", None)
                if image is None:
                    continue
                asset = image_service.upload_embedded(
                    image.blob,
                    extension=image.ext,
                    source=f"{self.context.file_name}:slide-{slide_index + 1}",
                )
                if asset:
                    image_map[(slide_index, shape.shape_id)] = asset.markdown
                    image_files.append(asset.file)
        return image_map, image_files

    @staticmethod
    def _table_to_markdown(table) -> str:
        if not table.rows:
            return ""
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        rows = [[cell.text.strip() for cell in row.cells] for row in list(table.rows)[1:]]
        return render_markdown_table(headers, rows)

    def _parse_pptx(self) -> tuple[str, list]:
        presentation = Presentation(BytesIO(self.context.file_bytes))
        image_map, image_files = self._extract_images(presentation)
        content: list[str] = []

        for slide_index, slide in enumerate(presentation.slides):
            slide_content: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        paragraph_parts: list[str] = []
                        for run in paragraph.runs:
                            run_text = run.text or ""
                            address = run.hyperlink.address if run.hyperlink else None
                            if address and self._is_safe_hyperlink(address):
                                paragraph_parts.append(f"[{run_text}]({address})")
                            else:
                                paragraph_parts.append(run_text)
                        paragraph_text = "".join(paragraph_parts).strip()
                        if paragraph_text:
                            slide_content.append(paragraph_text)

                image_markdown = image_map.get((slide_index, shape.shape_id))
                if image_markdown:
                    slide_content.append(image_markdown)

                if shape.has_table:
                    table_markdown = self._table_to_markdown(shape.table)
                    if table_markdown:
                        slide_content.append(table_markdown.rstrip())

            if slide_content:
                content.append(f"# Slide {slide_index + 1}\n" + "\n".join(slide_content))
        return "\n\n".join(content), image_files

    @staticmethod
    def _is_safe_hyperlink(url: str) -> bool:
        parsed = urlparse(url)
        if parsed.scheme.casefold() == "mailto":
            return bool(parsed.path)
        return parsed.scheme.casefold() in {"http", "https"} and bool(parsed.netloc)
